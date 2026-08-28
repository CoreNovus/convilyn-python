"""PPTX and XLSX extraction — the two biggest holes in the old format matrix.

Neither format could reach Markdown at all before this: ``SUPPORTED_CONVERSIONS``
never listed ``md`` as a target for either.
"""

from __future__ import annotations

import io
from datetime import datetime
from pathlib import Path

import pytest

from convilyn.local._engine.markdown.pptx import extract as extract_pptx
from convilyn.local._engine.markdown.xlsx import extract as extract_xlsx


def _png(size: tuple[int, int] = (256, 256)) -> bytes:
    from PIL import Image

    image = Image.new("RGB", size)
    image.putdata(
        [
            ((x * 7) % 256, (y * 11) % 256, (x * y) % 256)
            for y in range(size[1])
            for x in range(size[0])
        ]
    )
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


# ── PPTX ─────────────────────────────────────────────────────────────


@pytest.fixture
def deck(tmp_path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Results"
    slide.placeholders[1].text_frame.text = "First point"
    slide.notes_slide.notes_text_frame.text = "Say the number out loud."

    second = presentation.slides.add_slide(presentation.slide_layouts[5])
    second.shapes.title.text = "Appendix"
    image_path = tmp_path / "f.png"
    image_path.write_bytes(_png())
    second.shapes.add_picture(str(image_path), Inches(1), Inches(1))

    path = tmp_path / "d.pptx"
    presentation.save(str(path))
    return path


class TestPptx:
    def test_slide_title_becomes_a_heading(self, deck):
        headings = [b.text for b in extract_pptx(deck).blocks if b.kind == "heading"]

        assert "Results" in headings

    def test_body_placeholder_becomes_a_list_item(self, deck):
        items = [b.text for b in extract_pptx(deck).blocks if b.kind == "list_item"]

        assert "First point" in items

    def test_speaker_notes_are_kept_as_a_quote(self, deck):
        """Notes are usually the densest text in a deck; dropping them loses it."""
        quotes = [b.text for b in extract_pptx(deck).blocks if b.kind == "quote"]

        assert "Say the number out loud." in quotes

    def test_slides_are_separated_by_a_page_break(self, deck):
        assert any(b.kind == "page_break" for b in extract_pptx(deck).blocks)

    def test_picture_is_extracted_with_bytes(self, deck):
        image = next(b.image for b in extract_pptx(deck).blocks if b.kind == "image")

        assert image is not None and image.data.startswith(b"\x89PNG")

    def test_untitled_slide_still_gets_a_heading(self, tmp_path):
        """A blank layout has no title; the section must still be navigable."""
        from pptx import Presentation

        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        path = tmp_path / "b.pptx"
        presentation.save(str(path))

        assert extract_pptx(path).blocks[0].kind == "heading"

    def test_source_format_is_reported(self, deck):
        assert extract_pptx(deck).source_format == "pptx"


# ── XLSX ─────────────────────────────────────────────────────────────


@pytest.fixture
def workbook(tmp_path: Path) -> Path:
    import openpyxl

    book = openpyxl.Workbook()
    sheet = book.active
    sheet.title = "Orders"
    sheet.append(["item", "qty", "when"])
    sheet.append(["bolt", 4, datetime(2026, 8, 9)])
    sheet.append(["nut", 12.0, datetime(2026, 8, 10, 14, 30)])

    second = book.create_sheet("Notes")
    second.append(["free text"])

    path = tmp_path / "w.xlsx"
    book.save(str(path))
    return path


class TestXlsx:
    def test_each_sheet_becomes_a_heading(self, workbook):
        headings = [b.text for b in extract_xlsx(workbook).blocks if b.kind == "heading"]

        assert headings == ["Orders", "Notes"]

    def test_rows_become_a_table(self, workbook):
        table = next(b for b in extract_xlsx(workbook).blocks if b.kind == "table")

        assert table.rows[0] == ("item", "qty", "when")

    def test_whole_number_float_loses_its_decimal(self, workbook):
        """openpyxl returns 12.0 for an integer cell; "12.0" reads as wrong."""
        table = next(b for b in extract_xlsx(workbook).blocks if b.kind == "table")

        assert table.rows[2][1] == "12"

    def test_midnight_datetime_renders_as_a_plain_date(self, workbook):
        table = next(b for b in extract_xlsx(workbook).blocks if b.kind == "table")

        assert table.rows[1][2] == "2026-08-09"

    def test_datetime_with_a_time_keeps_it(self, workbook):
        table = next(b for b in extract_xlsx(workbook).blocks if b.kind == "table")

        assert table.rows[2][2] == "2026-08-10 14:30:00"

    def test_empty_sheet_is_reported_not_dropped(self, tmp_path):
        import openpyxl

        book = openpyxl.Workbook()
        book.active.title = "Blank"
        path = tmp_path / "e.xlsx"
        book.save(str(path))

        texts = [b.text for b in extract_xlsx(path).blocks if b.kind == "paragraph"]

        assert "(empty sheet)" in texts

    def test_source_format_is_reported(self, workbook):
        assert extract_xlsx(workbook).source_format == "xlsx"


class TestASingleSheetGetsNoHeading:
    """A heading exists to navigate BETWEEN sections; a single-sheet workbook
    has nothing to navigate between. Regression for the sheet name (most
    commonly Excel's own UI default, "Sheet1", never renamed) leaking into
    the output as if it were the document's own first line of content."""

    def test_the_sheet_name_does_not_become_a_heading(self, tmp_path):
        import openpyxl

        book = openpyxl.Workbook()
        book.active.title = "Sheet1"
        book.active.append(["item", "qty"])
        book.active.append(["bolt", 4])
        path = tmp_path / "single.xlsx"
        book.save(str(path))

        headings = [b.text for b in extract_xlsx(path).blocks if b.kind == "heading"]

        assert headings == []

    def test_the_omission_carries_no_warning(self, tmp_path):
        """Deliberately silent, not an oversight: a single-sheet workbook is
        the commonest shape this function sees, and warning on every one of
        them would be the exact cry-wolf pattern
        `workbook_guard.py::assert_no_sheets_would_be_dropped` already
        rejected for the same population (#4111) — "a workbook measured at
        exactly one sheet loses nothing, and warning anyway would cry wolf
        on the commonest workbook conversion here". A sheet named
        "CONFIDENTIAL" is the harder case for this rule, not an exception to
        it: the count is still the only signal, matching the module's own
        stated design (sheet count, not sheet name, decides).

        Unlike its sibling tests in this class, this one is a forward guard
        rather than a regression pin: the pre-fix code also emitted no
        warning here (it emitted a heading instead), so this assertion
        passes on both sides of the fix. Its job is to stop a FUTURE change
        from re-adding a warning for this omission, not to detect the
        original bug."""
        import openpyxl

        book = openpyxl.Workbook()
        book.active.title = "CONFIDENTIAL"
        book.active.append(["item", "qty"])
        path = tmp_path / "confidential.xlsx"
        book.save(str(path))

        assert extract_xlsx(path).warnings == ()

    def test_the_table_still_survives(self, tmp_path):
        import openpyxl

        book = openpyxl.Workbook()
        book.active.title = "Sheet1"
        book.active.append(["item", "qty"])
        book.active.append(["bolt", 4])
        path = tmp_path / "single_table.xlsx"
        book.save(str(path))

        table = next(b for b in extract_xlsx(path).blocks if b.kind == "table")

        assert table.rows[0] == ("item", "qty")

    def test_a_deliberately_named_single_sheet_also_gets_no_heading(self, tmp_path):
        """Not keyed on the name looking like a default — sheet COUNT is the
        signal, so even a meaningfully-named single sheet gets none, the
        same way a single-page PDF gets no invented section label."""
        import openpyxl

        book = openpyxl.Workbook()
        book.active.title = "Q3 Report"
        book.active.append(["item", "qty"])
        path = tmp_path / "named_single.xlsx"
        book.save(str(path))

        headings = [b.text for b in extract_xlsx(path).blocks if b.kind == "heading"]

        assert headings == []

    def test_two_sheets_still_both_get_headings(self, workbook):
        """The guard against over-fixing: this must not become 'never emit
        a sheet heading' — multiple sheets genuinely need labels to tell
        them apart, and the existing `workbook` fixture (2 sheets) already
        covers this, restated here to make the boundary explicit next to
        the single-sheet cases above."""
        headings = [b.text for b in extract_xlsx(workbook).blocks if b.kind == "heading"]

        assert headings == ["Orders", "Notes"]


# ── number formats, in the tree that ships ───────────────────────────
#
# The rule: a number format is applied when doing so is LOSSLESS and ADDS
# meaning, and refused when it only reshapes appearance. Pinned here as well as
# upstream because this module's docstring — hand-written, installed by the
# projection — is what a user reads, and it now states this rule. Nothing else
# checks that a projected doc tells the truth.


def _one_cell(tmp_path, value, number_format: str) -> str:
    """The rendered text of a single formatted cell."""
    import openpyxl

    book = openpyxl.Workbook()
    cell = book.active.cell(row=1, column=1, value=value)
    cell.number_format = number_format
    path = tmp_path / "one.xlsx"
    book.save(str(path))

    table = next(b for b in extract_xlsx(path).blocks if b.kind == "table")
    return table.rows[0][0]


class TestNumberFormats:
    def test_a_percentage_is_shifted_and_marked(self, tmp_path):
        assert _one_cell(tmp_path, 0.279613264457963, "0.00%") == "27.9613264457963%"

    def test_the_percentage_format_does_not_round(self, tmp_path):
        """A spreadsheet shows ``-72.0%``; that drops eleven digits the file has."""
        assert _one_cell(tmp_path, -0.720386735542037, "0.0%") == "-72.0386735542037%"

    def test_the_shift_is_exact_not_arithmetic(self, tmp_path):
        """``0.07 * 100`` is ``7.000000000000001``. Binary float multiplication
        is not a decimal shift, and inventing digits is the tampering this rule
        forbids — while passing any test that only looked for a ``%``."""
        assert _one_cell(tmp_path, 0.07, "0%") == "7%"

    def test_a_currency_symbol_survives(self, tmp_path):
        assert _one_cell(tmp_path, -141270, '"NT$"#,##0') == "NT$-141270"

    def test_thousands_separators_are_not_applied(self, tmp_path):
        """Lossless, but zero semantic, and they break downstream parsing."""
        assert _one_cell(tmp_path, 1234567, "#,##0") == "1234567"

    def test_decimal_rounding_is_never_applied(self, tmp_path):
        assert _one_cell(tmp_path, 3.14159265, "0.00") == "3.14159265"
